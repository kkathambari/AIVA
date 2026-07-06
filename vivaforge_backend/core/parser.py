import pdfplumber
import docx
import pptx
import os

class DocumentParser:
    @staticmethod
    def extract_text(file_path: str) -> str:
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            return DocumentParser._extract_pdf(file_path)
        elif ext == ".docx":
            return DocumentParser._extract_docx(file_path)
        elif ext == ".pptx":
            return DocumentParser._extract_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file format: {ext}")

    @staticmethod
    def _extract_pdf(file_path: str) -> str:
        text = ""
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
        return text

    @staticmethod
    def _extract_docx(file_path: str) -> str:
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])

    @staticmethod
    def _extract_pptx(file_path: str) -> str:
        prs = pptx.Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
